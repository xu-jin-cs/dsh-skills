"""general/parsers.py — 6 格式解析实现（ETLEngine 自实现，零历史平台代码依赖）。

统一入口：parse(file_path, ext, rules) -> dict。
rules 为 parsing.yaml 对应节内容，由调用方（general_parse 执行器）注入；
本文件禁写任何业务规则默认值——规则键缺失即 KeyError 报错（契约铁律 2）。
保护标记（TABLE/SLIDE/CODE 起止）亦从 rules["markers"] 注入，本文件不定义字面量。

rules 期望结构（键缺失即报错）：
  parser_versions: {text, word, spreadsheet, presentation, html, office_binary}
  markers: {table_start, table_end, slide_start, slide_end, code_start, code_end}
  markdown: {heading_levels, keep_code_blocks, drop_horizontal_rules}
  word: {heading_styles, fill_merged_cells}
  spreadsheet: {header_candidate_rows, empty_row_drop, max_rows_per_sheet, output_format}
  presentation: {keep_notes, drop_master_text}
  html: {node_blacklist, heading_tags, keep_img_alt, unescape_html}
  fallback: {csv_encoding_detect_order, binary_convert_tools, convert_timeout_seconds}
"""
import csv
import html as html_module
import logging
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from engine.general.filehash import md5_file

logger = logging.getLogger("etl_engine.general.parsers")


class CorruptFileError(Exception):
    """文件损坏或全部降级链失败（调用方按 retry_exception.yaml 分流为 fatal）。"""


def _need(node: dict, key: str) -> Any:
    """规则键缺失即报错：执行器层禁写默认值（契约铁律 2）。"""
    if key not in node:
        raise KeyError(f"parsing rules 缺键: {key}")
    return node[key]


def _result(sections: list[dict], parser_name: str, parser_version: str,
            is_degrade: bool, md5: str) -> dict:
    """归一化输出（对齐旧引擎结构，供下游 clean/chunk 消费）。"""
    return {
        "sections": sections,
        "raw_text": "\n".join(s["text"] for s in sections),
        "parser_name": parser_name,
        "parser_version": parser_version,
        "is_degrade": is_degrade,
        "md5": md5,
    }


def _wrap(text: str, start: str, end: str) -> str:
    return f"{start}\n{text}\n{end}"


# ---------------------------------------------------------------- txt / md

def _parse_txt(path: Path, rules: dict) -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="replace")
    sections = []
    for para in re.split(r"\n\s*\n", text.strip()):
        cleaned = para.replace("\n", " ").strip()
        if cleaned:
            sections.append({"type": "paragraph", "text": cleaned, "metadata": {}})
    return sections


def _parse_md(path: Path, rules: dict) -> list[dict]:
    md_rules = _need(rules, "markdown")
    markers = _need(rules, "markers")
    heading_levels = _need(md_rules, "heading_levels")
    keep_code = _need(md_rules, "keep_code_blocks")
    drop_hr = _need(md_rules, "drop_horizontal_rules")

    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()
    sections: list[dict] = []
    idx = 0
    while idx < len(lines):
        stripped = lines[idx].strip()
        if not stripped:
            idx += 1
            continue
        # 代码块
        if stripped.startswith("```"):
            lang = stripped[3:].strip()
            idx += 1
            buf = []
            while idx < len(lines) and not lines[idx].strip().startswith("```"):
                buf.append(lines[idx])
                idx += 1
            idx += 1  # 跳过收尾 fence（或文件结束）
            if keep_code and buf:
                sections.append({
                    "type": "code",
                    "text": _wrap("\n".join(buf), _need(markers, "code_start"),
                                  _need(markers, "code_end")),
                    "metadata": {"language": lang},
                })
            continue
        # ATX 标题（仅 heading_levels 内层级提为 title）
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m:
            level = len(m.group(1))
            title = m.group(2).strip()
            if level in heading_levels:
                sections.append({"type": "title", "text": title,
                                 "metadata": {"level": level, "marker": "#" * level}})
            elif title:
                sections.append({"type": "paragraph", "text": title,
                                 "metadata": {"level": level, "marker": "#" * level}})
            idx += 1
            continue
        # 水平线
        if drop_hr and re.match(r"^([-*_])\s*\1\s*\1+", stripped):
            idx += 1
            continue
        # 表格（GFM，行保留，保护标记包裹）
        if "|" in stripped and idx + 1 < len(lines) and "|" in lines[idx + 1]:
            rows = []
            while idx < len(lines) and "|" in lines[idx].strip() and lines[idx].strip():
                rows.append(lines[idx].strip())
                idx += 1
            sections.append({
                "type": "table",
                "text": _wrap("\n".join(rows), _need(markers, "table_start"),
                              _need(markers, "table_end")),
                "metadata": {"row_count": len(rows)},
            })
            continue
        # 列表（连续项并为单节）
        if re.match(r"^[-*+]\s+", stripped) or re.match(r"^\d+\.\s+", stripped):
            ordered = bool(re.match(r"^\d+\.\s+", stripped))
            buf = []
            while idx < len(lines):
                item = lines[idx].strip()
                if not item:
                    idx += 1
                    continue
                if not (re.match(r"^[-*+]\s+", item) or re.match(r"^\d+\.\s+", item)):
                    break
                buf.append(item)
                idx += 1
            if buf:
                sections.append({"type": "list", "text": "\n".join(buf),
                                 "metadata": {"ordered": ordered}})
            continue
        # 引用块 → 段落
        if stripped.startswith(">"):
            buf = []
            while idx < len(lines) and lines[idx].strip().startswith(">"):
                buf.append(lines[idx].strip()[1:].strip())
                idx += 1
            if buf:
                sections.append({"type": "paragraph", "text": " ".join(buf),
                                 "metadata": {"quote": True}})
            continue
        # 普通段落（聚合连续行）
        buf = []
        while idx < len(lines):
            line = lines[idx].strip()
            if not line or line.startswith("#") or line.startswith("```") or line.startswith(">"):
                break
            if re.match(r"^[-*+]\s+", line) or re.match(r"^\d+\.\s+", line):
                break
            if drop_hr and re.match(r"^([-*_])\s*\1\s*\1+", line):
                break
            if "|" in line and idx + 1 < len(lines) and "|" in lines[idx + 1]:
                break
            buf.append(line)
            idx += 1
        if buf:
            sections.append({"type": "paragraph", "text": " ".join(buf), "metadata": {}})
    return sections


# ---------------------------------------------------------------- csv / xlsx / xls

def _detect_header_rows(rows: list[list[str]], candidate: int) -> int:
    """启发式表头行数判定（0/1/2）：前两行均为纯文本无数字视为多级表头。"""
    if candidate <= 0 or not rows:
        return 0
    if candidate == 1 or len(rows) <= 1:
        return 1

    def _looks_like_header(r: list[str]) -> bool:
        return all(c and not any(ch.isdigit() for ch in c) for c in r)

    if _looks_like_header(rows[0]) and _looks_like_header(rows[1]):
        return 2
    return 1


def _format_sheet(rows: list[list[str]], rules: dict) -> str:
    sp = _need(rules, "spreadsheet")
    markers = _need(rules, "markers")
    header_rows = _detect_header_rows(rows, _need(sp, "header_candidate_rows"))
    header = rows[0] if header_rows >= 1 else []
    if header_rows > 1:  # 多级表头合并为单列名
        width = max(len(r) for r in rows[:header_rows])
        header = ["_".join(p for p in (r[i] if i < len(r) else "" for r in rows[:header_rows]) if p)
                  for i in range(width)]
    lines = []
    for row in rows[header_rows:]:
        if _need(sp, "output_format") == "key_value" and header:
            pairs = [f"{header[i] if i < len(header) else 'col_' + str(i)}:{val}"
                     for i, val in enumerate(row)]
            lines.append(", ".join(pairs))
        else:
            lines.append(", ".join(row))
    return _wrap("\n".join(lines), _need(markers, "table_start"),
                 _need(markers, "table_end"))


def _parse_csv(path: Path, rules: dict) -> list[dict]:
    fb = _need(rules, "fallback")
    sp = _need(rules, "spreadsheet")
    text = None
    used = ""
    for enc in _need(fb, "csv_encoding_detect_order"):
        try:
            text = path.read_text(encoding=enc)
            used = enc
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if text is None:
        raise CorruptFileError(f"csv 编码探测全部失败: {path.name}")
    from io import StringIO
    rows = [[cell.strip() for cell in row] for row in csv.reader(StringIO(text))]
    if _need(sp, "empty_row_drop"):
        rows = [r for r in rows if any(r)]
    if not rows:
        return []
    return [{"type": "table", "text": _format_sheet(rows, rules),
             "metadata": {"sheet_name": path.stem, "row_count": len(rows),
                          "encoding": used}}]


def _parse_xlsx(path: Path, rules: dict, ext: str) -> tuple[list[dict], bool]:
    """xlsx 原生解析；xls 走 openpyxl 尝试，失败降级（is_degrade=True）。"""
    sp = _need(rules, "spreadsheet")
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    except Exception as exc:
        if ext == "xls":
            logger.warning("xls 降级（openpyxl 不支持 .xls）: %s", exc)
            return [], True
        raise CorruptFileError(f"xlsx 解析失败: {exc}") from exc
    max_rows = _need(sp, "max_rows_per_sheet")
    sections = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            truncated = False
            for ridx, row in enumerate(ws.iter_rows(values_only=True)):
                if ridx >= max_rows:
                    truncated = True
                    break
                rows.append([str(c) if c is not None else "" for c in row])
            if _need(sp, "empty_row_drop"):
                rows = [r for r in rows if any(r)]
            if not rows:
                continue
            sections.append({
                "type": "table",
                "text": _format_sheet(rows, rules),
                "metadata": {"sheet_name": sheet_name, "row_count": len(rows),
                             "truncated": truncated},
            })
    finally:
        wb.close()
    return sections, False


# ---------------------------------------------------------------- docx

def _parse_docx(path: Path, rules: dict) -> list[dict]:
    word = _need(rules, "word")
    markers = _need(rules, "markers")
    heading_styles = _need(word, "heading_styles")
    fill_merged = _need(word, "fill_merged_cells")
    try:
        import docx
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise CorruptFileError(f"python-docx 不可用: {exc}") from exc
    try:
        doc = docx.Document(str(path))
    except Exception as exc:
        raise CorruptFileError(f"docx 打开失败: {exc}") from exc

    def _heading_level(style_name: str) -> int:
        if style_name in heading_styles:
            return heading_styles.index(style_name) + 1
        if style_name.startswith("Heading "):
            try:
                return int(style_name.split(" ", 1)[1])
            except ValueError:
                pass
        return 0

    def _table_section(table: Any) -> dict | None:
        lines = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if fill_merged:  # 横向前向填充合并单元格
                for i in range(1, len(cells)):
                    if not cells[i]:
                        cells[i] = cells[i - 1]
            lines.append("\t".join(cells))
        if not lines:
            return None
        return {"type": "table",
                "text": _wrap("\n".join(lines), _need(markers, "table_start"),
                              _need(markers, "table_end")),
                "metadata": {"row_count": len(table.rows)}}

    # 按文档正文顺序交错提取段落与表格
    sections = []
    from docx.oxml.ns import qn
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            level = _heading_level(style)
            sections.append({"type": "title" if level > 0 else "paragraph",
                             "text": text,
                             "metadata": {"style": style, "level": level}})
        elif child.tag == qn("w:tbl"):
            sec = _table_section(Table(child, doc))
            if sec:
                sections.append(sec)
    return sections


# ---------------------------------------------------------------- pptx

def _parse_pptx(path: Path, rules: dict) -> list[dict]:
    pres = _need(rules, "presentation")
    markers = _need(rules, "markers")
    keep_notes = _need(pres, "keep_notes")
    drop_master = _need(pres, "drop_master_text")
    try:
        import pptx
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError as exc:
        raise CorruptFileError(f"python-pptx 不可用: {exc}") from exc
    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
    try:
        prs = pptx.Presentation(str(path))
    except Exception as exc:
        raise CorruptFileError(f"pptx 打开失败: {exc}") from exc

    sections = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        master_texts: set[str] = set()
        if drop_master:  # 收集版式/母版占位文本，正文命中即剔除
            for holder in (slide.slide_layout, slide.slide_layout.slide_master):
                for shape in holder.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        master_texts.add(shape.text.strip())
        title = ""
        body = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text or text in master_texts:
                continue
            if shape.is_placeholder and shape.placeholder_format.type in title_types:
                title = text
                continue
            body.append(text)
        parts = ([title] if title else []) + body
        if keep_notes and slide.has_notes_slide and slide.notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(note)
        if not parts:
            continue
        sections.append({
            "type": "slide",
            "text": _wrap("\n".join(parts), _need(markers, "slide_start"),
                          _need(markers, "slide_end")),
            "metadata": {"slide_no": slide_no, "slide_title": title},
        })
    return sections


# ---------------------------------------------------------------- html / htm

def _parse_html(path: Path, rules: dict) -> list[dict]:
    html_rules = _need(rules, "html")
    markers = _need(rules, "markers")
    from bs4 import BeautifulSoup

    text = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(text, "html.parser")
    for tag in _need(html_rules, "node_blacklist"):
        for node in soup.find_all(tag):
            node.decompose()
    if _need(html_rules, "keep_img_alt"):
        for img in soup.find_all("img"):
            alt = img.get("alt", "")
            if alt:
                img.replace_with(alt)
    else:
        for img in soup.find_all("img"):
            img.decompose()

    heading_tags = _need(html_rules, "heading_tags")
    unescape = _need(html_rules, "unescape_html")

    def _clean(t: str) -> str:
        t = re.sub(r"\s+", " ", t).strip()
        return html_module.unescape(t) if unescape else t

    sections = []
    container = soup.find("body") or soup
    for elem in container.descendants:
        if not getattr(elem, "name", None):
            continue
        if elem.name in heading_tags:
            txt = _clean(elem.get_text(" ", strip=True))
            if txt:
                sections.append({"type": "title", "text": txt,
                                 "metadata": {"tag": elem.name,
                                              "level": int(elem.name[1])}})
        elif elem.name == "pre":
            code = elem.get_text("\n", strip=True)
            if code:
                sections.append({"type": "code",
                                 "text": _wrap(code, _need(markers, "code_start"),
                                               _need(markers, "code_end")),
                                 "metadata": {}})
        elif elem.name == "table":
            rows = ["\t".join(td.get_text(" ", strip=True)
                              for td in tr.find_all(["td", "th"]))
                    for tr in elem.find_all("tr")]
            rows = [r for r in rows if r.strip()]
            if rows:
                sections.append({"type": "table",
                                 "text": _wrap("\n".join(rows),
                                               _need(markers, "table_start"),
                                               _need(markers, "table_end")),
                                 "metadata": {"row_count": len(rows)}})
        elif elem.name in ("p",):
            txt = _clean(elem.get_text(" ", strip=True))
            if txt and not re.match(r"^[\s\|\-\*\_\.]+$", txt):
                sections.append({"type": "paragraph", "text": txt, "metadata": {}})
    # 去重（嵌套标签可能重复产出同文本）
    seen, unique = set(), []
    for s in sections:
        if s["text"] not in seen:
            seen.add(s["text"])
            unique.append(s)
    return unique


# ---------------------------------------------------------------- doc / wps 降级链

def _convert_with_tool(tool: str, path: Path, timeout: int) -> str:
    """按工具名执行二进制→txt 转换（工具清单由规则注入，此处仅行为实现）。"""
    if tool == "textutil":
        proc = subprocess.run(["textutil", "-convert", "txt", "-stdout", str(path)],
                              capture_output=True, text=True, timeout=timeout, check=False)
        return proc.stdout if proc.returncode == 0 else ""
    if tool == "libreoffice":
        if not shutil.which("libreoffice"):
            return ""
        tmpdir = Path(tempfile.mkdtemp(prefix="general_etl_bin_"))
        try:
            subprocess.run(["libreoffice", "--headless", "--convert-to", "txt:Text",
                            "--outdir", str(tmpdir), str(path)],
                           capture_output=True, text=True, timeout=timeout, check=False)
            out = tmpdir / (path.stem + ".txt")
            return out.read_text(encoding="utf-8", errors="replace") if out.exists() else ""
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)
    logger.warning("未识别的降级工具（跳过）: %s", tool)
    return ""


def _parse_office_binary(path: Path, rules: dict) -> list[dict]:
    fb = _need(rules, "fallback")
    timeout = _need(fb, "convert_timeout_seconds")
    for tool in _need(fb, "binary_convert_tools"):
        try:
            text = _convert_with_tool(tool, path, timeout)
        except Exception as exc:
            logger.warning("降级工具 %s 失败: %s", tool, exc)
            continue
        if text.strip():
            return [{"type": "paragraph", "text": line.strip(),
                     "metadata": {"degraded": True, "tool": tool}}
                    for line in text.split("\n") if line.strip()]
    raise CorruptFileError(f"二进制 Office 降级链全部失败: {path.name}")


# ---------------------------------------------------------------- 统一入口

def parse(file_path: str, ext: str, rules: dict) -> dict:
    """6 格式统一解析入口。rules 由调用方从 parsing.yaml 注入（缺键即报错）。

    返回归一化结构：{sections, raw_text, parser_name, parser_version, is_degrade, md5}
    """
    path = Path(file_path)
    ext = ext.lower().lstrip(".")
    versions = _need(rules, "parser_versions")
    md5 = md5_file(path)

    if ext == "txt":
        return _result(_parse_txt(path, rules), "text_parser",
                       str(_need(versions, "text")), False, md5)
    if ext == "json":
        # 2026-08-17 用户裁定白名单扩列：json 按纯文本解析，复用 text_parser
        return _result(_parse_txt(path, rules), "text_parser",
                       str(_need(versions, "text")), False, md5)
    if ext == "md":
        return _result(_parse_md(path, rules), "text_parser",
                       str(_need(versions, "text")), False, md5)
    if ext == "csv":
        return _result(_parse_csv(path, rules), "spreadsheet_parser",
                       str(_need(versions, "spreadsheet")), False, md5)
    if ext in ("xlsx", "xls"):
        sections, degraded = _parse_xlsx(path, rules, ext)
        return _result(sections, "spreadsheet_parser",
                       str(_need(versions, "spreadsheet")), degraded, md5)
    if ext == "docx":
        return _result(_parse_docx(path, rules), "docx_parser",
                       str(_need(versions, "word")), False, md5)
    if ext in ("doc", "wps"):
        return _result(_parse_office_binary(path, rules), "office_binary_parser",
                       str(_need(versions, "office_binary")), True, md5)
    if ext == "pptx":
        return _result(_parse_pptx(path, rules), "presentation_parser",
                       str(_need(versions, "presentation")), False, md5)
    if ext == "ppt":
        # 二进制 ppt 走外部工具降级链（U3 集成补登：whitelist 含 ppt 但入口未路由）
        return _result(_parse_office_binary(path, rules), "office_binary_parser",
                       str(_need(versions, "office_binary")), True, md5)
    if ext in ("html", "htm"):
        return _result(_parse_html(path, rules), "html_parser",
                       str(_need(versions, "html")), False, md5)
    raise ValueError(f"未登记的解析格式（应由白名单层拦截）: {ext}")
